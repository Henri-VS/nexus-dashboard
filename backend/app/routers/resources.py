import base64
import io

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel

router = APIRouter()

# ── ChromaDB (lazy, ephemeral) ────────────────────────────────────────────────

_collection = None


def _get_collection():
    global _collection
    if _collection is not None:
        return _collection
    try:
        import chromadb  # type: ignore
    except ImportError:
        raise HTTPException(status_code=503, detail="chromadb is not installed. Run: pip install chromadb")
    client = chromadb.Client()
    _collection = client.get_or_create_collection("resources")
    return _collection


# ── Text extraction ───────────────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        raise HTTPException(status_code=503, detail="pypdf is not installed. Run: pip install pypdf")
    reader = PdfReader(io.BytesIO(data))
    pages = [p.extract_text() or "" for p in reader.pages]
    return "\n\n".join(p for p in pages if p.strip())


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise HTTPException(status_code=503, detail="python-pptx is not installed. Run: pip install python-pptx")
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
    return "\n\n".join(slides)


# ── Chunking ──────────────────────────────────────────────────────────────────

def _chunk(text: str, size: int = 500, overlap: int = 50) -> list[str]:
    words = text.split()
    if not words:
        return []
    chunks: list[str] = []
    i = 0
    while i < len(words):
        chunks.append(" ".join(words[i : i + size]))
        i += size - overlap
    return chunks


# ── Endpoints ─────────────────────────────────────────────────────────────────

class ProcessRequest(BaseModel):
    filename: str
    content: str          # base64-encoded file bytes
    file_type: str        # 'pdf' | 'pptx' | 'ppt' | 'md'
    subject: str = ""
    semester: str = ""


@router.post("/process")
async def process_resource(req: ProcessRequest):
    """
    Extract text from a resource file, chunk it, and store it in ChromaDB.
    Returns the extracted text (truncated to 50 000 chars) for the frontend to cache.
    """
    try:
        raw = base64.b64decode(req.content)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid base64 content")

    ft = req.file_type.lower().lstrip(".")

    if ft == "pdf":
        text = _extract_pdf(raw)
    elif ft in ("pptx", "ppt"):
        text = _extract_pptx(raw)
    elif ft == "md":
        text = raw.decode("utf-8", errors="replace")
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ft}")

    if not text.strip():
        raise HTTPException(status_code=422, detail="No text could be extracted from the file")

    chunks = _chunk(text)
    if not chunks:
        raise HTTPException(status_code=422, detail="No content after chunking")

    coll = _get_collection()

    # Remove any existing chunks for this filename
    try:
        existing = coll.get(where={"filename": req.filename})
        if existing.get("ids"):
            coll.delete(ids=existing["ids"])
    except Exception:
        pass

    ids = [f"{req.filename}::{i}" for i in range(len(chunks))]
    coll.add(
        documents=chunks,
        ids=ids,
        metadatas=[
            {
                "filename": req.filename,
                "subject": req.subject,
                "semester": req.semester,
                "chunk_index": str(i),
            }
            for i in range(len(chunks))
        ],
    )

    return {
        "status": "processed",
        "chunks": len(chunks),
        "filename": req.filename,
        "text": text[:50_000],
    }


@router.get("/search")
async def search_resources(q: str, subject: str = "", n: int = 5):
    """Query ChromaDB for chunks relevant to q, optionally filtered by subject."""
    if not q.strip():
        raise HTTPException(status_code=400, detail="Query cannot be empty")

    coll = _get_collection()

    # Guard: ChromaDB raises if collection is empty
    try:
        count = coll.count()
    except Exception:
        count = 0

    if count == 0:
        return {"query": q, "chunks": []}

    kwargs: dict = {"query_texts": [q], "n_results": min(n, count)}
    if subject:
        kwargs["where"] = {"subject": subject}

    try:
        results = coll.query(**kwargs)
    except Exception as exc:
        # If filtered query fails (e.g. subject not found), retry without filter
        try:
            kwargs.pop("where", None)
            results = coll.query(**kwargs)
        except Exception:
            raise HTTPException(status_code=500, detail=str(exc))

    docs = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    dists = results.get("distances", [[]])[0]

    return {
        "query": q,
        "chunks": [
            {
                "content": doc,
                "metadata": meta,
                "score": round(1.0 - dist, 4),
            }
            for doc, meta, dist in zip(docs, metas, dists)
        ],
    }
